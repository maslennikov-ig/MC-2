"""Generate deterministic quality fixtures for the Docling A/B benchmark."""

from pathlib import Path

from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm as DocxCm
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm, Pt
from reportlab.lib.pagesizes import A4
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas


ROOT = Path(__file__).resolve().parent

# Any Cyrillic-capable TTF. The first existing candidate wins so the generator
# runs on both the Windows and the WSL/Linux checkout without editing.
FONT_CANDIDATES = (
    Path("C:/Windows/Fonts/arial.ttf"),
    Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
    Path("/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"),
)


def _font_path() -> Path:
    for candidate in FONT_CANDIDATES:
        if candidate.exists():
            return candidate
    raise FileNotFoundError(
        "No Cyrillic-capable TTF found; add one to FONT_CANDIDATES: "
        + ", ".join(str(path) for path in FONT_CANDIDATES)
    )


def _append_value(parent, tag: str, value: str):
    element = OxmlElement(tag)
    element.set(qn("w:val"), value)
    parent.append(element)
    return element


def _add_multilevel_decimal_numbering(document: Document) -> int:
    numbering = document.part.numbering_part.element
    abstract_ids = [
        int(element.get(qn("w:abstractNumId")))
        for element in numbering.findall(qn("w:abstractNum"))
    ]
    num_ids = [
        int(element.get(qn("w:numId")))
        for element in numbering.findall(qn("w:num"))
    ]
    abstract_id = max(abstract_ids, default=-1) + 1
    num_id = max(num_ids, default=0) + 1

    abstract = OxmlElement("w:abstractNum")
    abstract.set(qn("w:abstractNumId"), str(abstract_id))
    _append_value(abstract, "w:multiLevelType", "multilevel")
    for level, marker, left in ((0, "%1.", 720), (1, "%1.%2.", 1440)):
        level_element = OxmlElement("w:lvl")
        level_element.set(qn("w:ilvl"), str(level))
        _append_value(level_element, "w:start", "1")
        _append_value(level_element, "w:numFmt", "decimal")
        _append_value(level_element, "w:lvlText", marker)
        _append_value(level_element, "w:lvlJc", "left")
        paragraph_properties = OxmlElement("w:pPr")
        indentation = OxmlElement("w:ind")
        indentation.set(qn("w:left"), str(left))
        indentation.set(qn("w:hanging"), "360")
        paragraph_properties.append(indentation)
        level_element.append(paragraph_properties)
        abstract.append(level_element)
    numbering.append(abstract)

    concrete = OxmlElement("w:num")
    concrete.set(qn("w:numId"), str(num_id))
    _append_value(concrete, "w:abstractNumId", str(abstract_id))
    numbering.append(concrete)
    return num_id


def _set_list_level(paragraph, num_id: int, level: int) -> None:
    num_properties = paragraph._p.get_or_add_pPr().get_or_add_numPr()
    indentation_level = num_properties.find(qn("w:ilvl"))
    if indentation_level is None:
        indentation_level = OxmlElement("w:ilvl")
        num_properties.insert(0, indentation_level)
    indentation_level.set(qn("w:val"), str(level))
    numbering_id = num_properties.find(qn("w:numId"))
    if numbering_id is None:
        numbering_id = OxmlElement("w:numId")
        num_properties.append(numbering_id)
    numbering_id.set(qn("w:val"), str(num_id))


def make_docx() -> None:
    document = Document()
    document.add_heading("Проверка структуры документа", level=1)
    document.add_heading("Нумерованные и вложенные списки", level=2)
    list_num_id = _add_multilevel_decimal_numbering(document)
    list_items = (
        ("Первый обязательный шаг", "List Number", 0),
        ("Вложенный шаг А", "List Number 2", 1),
        ("Вложенный шаг Б", "List Number 2", 1),
        ("Второй обязательный шаг", "List Number", 0),
    )
    for text, style, level in list_items:
        paragraph = document.add_paragraph(text, style=style)
        _set_list_level(paragraph, list_num_id, level)

    document.add_heading("Таблица с объединёнными ячейками", level=2)
    table = document.add_table(rows=3, cols=3)
    table.style = "Table Grid"
    merged = table.cell(0, 0).merge(table.cell(0, 2))
    merged.text = "Сводные показатели"
    merged.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
    rows = [
        ("Метрика", "Значение", "Единица"),
        ("Точность", "98", "процентов"),
    ]
    for row_index, values in enumerate(rows, start=1):
        for column_index, value in enumerate(values):
            table.cell(row_index, column_index).text = value
    for section in document.sections:
        section.top_margin = DocxCm(2)
        section.bottom_margin = DocxCm(2)
    document.save(ROOT / "structured-lists-and-table.docx")


def make_hierarchy_docx() -> None:
    """Three real heading levels, with values that are ambiguous without them.

    Both subsections state a "Целевое значение", and only the heading path says
    which product each belongs to. A chunking strategy that loses headings can
    still retrieve the number and cannot answer the question.
    """
    document = Document()
    document.add_heading("Регламент обработки документов", level=1)
    document.add_paragraph(
        "Документ описывает требования к двум продуктовым линиям и порядок их проверки."
    )

    document.add_heading("Линия Альфа", level=2)
    document.add_paragraph("Линия Альфа обслуживает корпоративных клиентов.")
    document.add_heading("Целевые показатели Альфы", level=3)
    document.add_paragraph("Целевое значение: 42 единицы в сутки при полной загрузке.")
    document.add_heading("Ограничения Альфы", level=3)
    document.add_paragraph("Максимальная задержка обработки составляет 15 минут.")

    document.add_heading("Линия Бета", level=2)
    document.add_paragraph("Линия Бета обслуживает розничных клиентов.")
    document.add_heading("Целевые показатели Беты", level=3)
    document.add_paragraph("Целевое значение: 77 единиц в сутки при полной загрузке.")
    document.add_heading("Ограничения Беты", level=3)
    document.add_paragraph("Максимальная задержка обработки составляет 30 минут.")

    document.save(ROOT / "hierarchy-multilevel.docx")


# Markers carry the hierarchy: Docling's numbering inference recognises `1.` as
# an arabic marker of depth 1 and `1.1.` as a dotted marker of depth 2, and the
# two distinct families are what it compresses into two heading levels.
NUMBERED_SECTIONS = (
    (1, "1. Введение", "Введение объясняет цель регламента и область его применения."),
    (2, "1.1. Область применения", "Регламент применяется ко всем входящим документам."),
    (2, "1.2. Термины", "Термин «партия» означает набор документов одной загрузки."),
    (1, "2. Порядок обработки", "Обработка выполняется в три последовательных этапа."),
    (2, "2.1. Приём", "На приёме проверяется формат и целостность файла."),
    (2, "2.2. Проверка", "Контрольное значение проверки равно 64 пунктам."),
    (1, "3. Заключение", "Заключение фиксирует ответственных и сроки пересмотра."),
)


def make_numbered_sections_pdf() -> None:
    """A text-layer PDF whose outline numbering encodes a two-level hierarchy.

    Docling reports every section header at level 1 unless heading-hierarchy
    inference is enabled, so this fixture is what proves the feature flag does
    something, and proves it on structure rather than on heading count.
    """
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    font_path = _font_path()
    bold_path = font_path.with_name(font_path.name.replace(".ttf", "-Bold.ttf"))
    pdfmetrics.registerFont(TTFont("FixtureSans", str(font_path)))
    pdfmetrics.registerFont(
        TTFont("FixtureSans-Bold", str(bold_path if bold_path.exists() else font_path))
    )

    pdf = canvas.Canvas(
        str(ROOT / "numbered-sections.pdf"), pagesize=A4, pageCompression=1
    )
    page_width, page_height = A4
    y = page_height - 80

    # Headings are bold and clearly larger than body text: the layout model has
    # to classify them as section headers before numbering can order them.
    pdf.setFont("FixtureSans-Bold", 22)
    pdf.drawString(70, y, "Регламент обработки: нумерованные разделы")
    y -= 60

    for level, heading, body in NUMBERED_SECTIONS:
        if y < 140:
            pdf.showPage()
            y = page_height - 80
        pdf.setFont("FixtureSans-Bold", 18 if level == 1 else 15)
        pdf.drawString(70, y, heading)
        y -= 30
        pdf.setFont("FixtureSans", 11)
        pdf.drawString(70, y, body)
        y -= 46

    pdf.showPage()
    pdf.save()


def add_textbox(slide, left, top, width, height, text, size=20, bold=False):
    shape = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    return shape


def make_pptx() -> None:
    presentation = Presentation()
    presentation.slide_width = Cm(33.867)
    presentation.slide_height = Cm(19.05)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    add_textbox(slide, 1.5, 0.7, 28, 1.4, "Порядок чтения и данные диаграммы", 28, True)
    add_textbox(slide, 2, 3.0, 8, 2, "Шаг один: собрать данные", 20, True)
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Cm(10.5), Cm(3.2), Cm(4), Cm(1.4))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(47, 111, 191)
    add_textbox(slide, 15, 3.0, 9, 2, "Шаг два: проверить результат", 20, True)
    add_textbox(slide, 2, 5.3, 22, 1, "Подпись схемы: последовательность обработки", 15)

    chart_data = ChartData()
    chart_data.categories = ["Квартал 1", "Квартал 2", "Квартал 3"]
    chart_data.add_series("Документы", (10, 20, 30))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Cm(2),
        Cm(7),
        Cm(22),
        Cm(9),
        chart_data,
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Обработанные документы по кварталам"
    add_textbox(slide, 24.8, 8, 7, 4, "Контрольные значения:\n10, 20, 30", 18, True)
    presentation.save(ROOT / "reading-order-chart.pptx")


def make_raster_ocr_pdf() -> None:
    width, height = 2480, 3508
    font_path = _font_path()
    image = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(image)
    font = ImageFont.truetype(str(font_path), 62)
    heading = ImageFont.truetype(str(font_path), 74)
    small = ImageFont.truetype(str(font_path), 48)

    draw.text((160, 180), "РУССКИЙ OCR: КОНТРОЛЬНЫЙ ДОКУМЕНТ", font=heading, fill="black")
    draw.text(
        (160, 390),
        "КОНТРОЛЬНАЯ ФРАЗА: КАЧЕСТВО РАСПОЗНАВАНИЯ",
        font=font,
        fill="black",
    )
    draw.text(
        (160, 520),
        "ВТОРАЯ ФРАЗА: ДОКУМЕНТЫ ПРЕВРАЩАЮТСЯ В ЗНАНИЯ",
        font=font,
        fill="black",
    )

    x_positions = [160, 1050, 1650, 2260]
    y_positions = [850, 1040, 1230, 1420]
    for x in x_positions:
        draw.line((x, y_positions[0], x, y_positions[-1]), fill="black", width=7)
    for y in y_positions:
        draw.line((x_positions[0], y, x_positions[-1], y), fill="black", width=7)
    cells = [
        ("Показатель", "Значение", "Единица"),
        ("Точность", "98", "процентов"),
        ("Страницы", "12", "штук"),
    ]
    for row, values in enumerate(cells):
        for column, value in enumerate(values):
            draw.text((x_positions[column] + 25, y_positions[row] + 45), value, font=small, fill="black")

    png_path = ROOT / ".raster-ocr-source.png"
    image.save(png_path, dpi=(300, 300))
    page_width, page_height = A4
    pdf = canvas.Canvas(str(ROOT / "russian-raster-ocr.pdf"), pagesize=A4, pageCompression=1)
    pdf.drawImage(ImageReader(str(png_path)), 0, 0, page_width, page_height)
    pdf.showPage()
    pdf.save()
    png_path.unlink()


PIXEL_FONT = {
    "N": ("10001", "11001", "10101", "10011", "10001", "10001", "10001"),
    "O": ("01110", "10001", "10001", "10001", "10001", "10001", "01110"),
    "T": ("11111", "00100", "00100", "00100", "00100", "00100", "00100"),
    "E": ("11111", "10000", "10000", "11110", "10000", "10000", "11111"),
    "X": ("10001", "10001", "01010", "00100", "01010", "10001", "10001"),
}


def make_vector_outline_pdf() -> None:
    pdf = canvas.Canvas(str(ROOT / "vector-outlines-no-text.pdf"), pagesize=A4, pageCompression=1)
    scale = 10
    x, y = 90, 500
    for character in "NO TEXT":
        if character == " ":
            x += scale * 4
            continue
        for row, pattern in enumerate(PIXEL_FONT[character]):
            for column, bit in enumerate(pattern):
                if bit == "1":
                    pdf.rect(x + column * scale, y - row * scale, scale, scale, stroke=0, fill=1)
        x += scale * 7
    # Vector-only diagram elements make the negative case non-empty visually.
    pdf.setLineWidth(3)
    pdf.circle(180, 300, 55, stroke=1, fill=0)
    pdf.line(235, 300, 360, 300)
    pdf.rect(360, 245, 110, 110, stroke=1, fill=0)
    pdf.showPage()
    pdf.save()


CODE_LISTING = (
    "def normalize_weight(raw_weight, minimum=0.0, maximum=1.0):",
    '    """Clamp a document weight into the accepted range."""',
    "    if raw_weight is None:",
    "        raise ValueError('weight is required')",
    "    return max(minimum, min(maximum, float(raw_weight)))",
)

# Values are deliberately unequal, three digits apart and not round, so a model
# that invents plausible numbers cannot land on them by accident.
CHART_SERIES = (("Альфа", 12), ("Бета", 34), ("Гамма", 56))


def _draw_formula(path: Path) -> None:
    """A quadratic formula as PIXELS, not as a text layer.

    Formula enrichment has to read the image; if the characters were selectable
    text the fixture would prove the PDF backend, not the model.
    """
    font = ImageFont.truetype(str(_font_path()), 44)
    small = ImageFont.truetype(str(_font_path()), 26)
    image = Image.new("RGB", (760, 190), "white")
    draw = ImageDraw.Draw(image)
    draw.text((30, 60), "x =", font=font, fill="black")
    draw.text((130, 25), "-b ±", font=font, fill="black")
    draw.text((250, 20), "√(b", font=font, fill="black")
    draw.text((330, 8), "2", font=small, fill="black")
    draw.text((355, 20), "- 4ac)", font=font, fill="black")
    draw.line((130, 88, 560, 88), fill="black", width=4)
    draw.text((300, 100), "2a", font=font, fill="black")
    image.save(path, format="PNG")


def _draw_bar_chart(path: Path) -> None:
    """A bar chart as PIXELS: no embedded chart XML to read the values from.

    The PPTX fixture already covers the native-chart path, where Docling
    recovers series straight from the source file. This one is the case that
    genuinely needs a vision model.
    """
    font = ImageFont.truetype(str(_font_path()), 22)
    title_font = ImageFont.truetype(str(_font_path()), 26)
    width, height = 720, 480
    image = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(image)
    draw.text((150, 20), "Обработано документов", font=title_font, fill="black")

    baseline = height - 70
    draw.line((90, baseline, width - 40, baseline), fill="black", width=3)
    draw.line((90, 70, 90, baseline), fill="black", width=3)

    scale = (baseline - 110) / max(value for _, value in CHART_SERIES)
    for index, (label, value) in enumerate(CHART_SERIES):
        left = 150 + index * 170
        top = baseline - value * scale
        draw.rectangle((left, top, left + 110, baseline), fill="#4c6ef5")
        draw.text((left + 30, top - 32), str(value), font=font, fill="black")
        draw.text((left + 20, baseline + 14), label, font=font, fill="black")

    image.save(path, format="PNG")


def make_enrichment_pdf() -> None:
    """One PDF carrying all three enrichment signals, each in its own section.

    Code stays a real text layer in a monospace font, because that is the shape
    the layout model classifies as code. The formula and the chart are raster
    images on purpose: they are the cases the advanced pass exists for.
    """
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    font_path = _font_path()
    bold_path = font_path.with_name(font_path.name.replace(".ttf", "-Bold.ttf"))
    pdfmetrics.registerFont(TTFont("FixtureSans", str(font_path)))
    pdfmetrics.registerFont(
        TTFont("FixtureSans-Bold", str(bold_path if bold_path.exists() else font_path))
    )

    formula_path = ROOT / "_enrichment-formula.png"
    chart_path = ROOT / "_enrichment-chart.png"
    _draw_formula(formula_path)
    _draw_bar_chart(chart_path)

    pdf = canvas.Canvas(
        str(ROOT / "enrichment-code-formula-chart.pdf"), pagesize=A4, pageCompression=1
    )
    page_width, page_height = A4
    y = page_height - 80

    pdf.setFont("FixtureSans-Bold", 20)
    pdf.drawString(60, y, "Справочник обработки документов")
    y -= 46

    pdf.setFont("FixtureSans-Bold", 15)
    pdf.drawString(60, y, "1. Листинг нормализации веса")
    y -= 30
    pdf.setFont("Courier", 10.5)
    for line in CODE_LISTING:
        pdf.drawString(72, y, line)
        y -= 16
    y -= 30

    pdf.setFont("FixtureSans-Bold", 15)
    pdf.drawString(60, y, "2. Формула корней квадратного уравнения")
    y -= 24
    pdf.drawImage(
        ImageReader(str(formula_path)), 60, y - 100, width=340, height=95, mask=None
    )
    y -= 140

    pdf.setFont("FixtureSans-Bold", 15)
    pdf.drawString(60, y, "3. Диаграмма обработанных документов")
    y -= 24
    pdf.drawImage(
        ImageReader(str(chart_path)), 60, y - 250, width=360, height=240, mask=None
    )

    pdf.showPage()
    pdf.save()

    formula_path.unlink()
    chart_path.unlink()


# ---------------------------------------------------------------------------
# Premium format families (Stage C)
#
# One fixture per family, each built around the structure that family exists to
# express, so a regression shows up as a missing FACT rather than as a diff in
# how much text survived. Every fixture is Russian, because that is what the
# platform ingests and because Cyrillic is where encoding regressions surface.
# ---------------------------------------------------------------------------


def _patch_xlsx_cached_formula(path: Path, cell: str, value: str) -> None:
    """Give a formula cell the cached result a real spreadsheet editor stores.

    MEASURED 2026-08-07: Docling reads the CACHED VALUE of a formula cell and
    never the expression. openpyxl writes `<f>SUM(B2:B3)</f><v></v>` because it
    does not evaluate anything, so a fixture saved straight from openpyxl proves
    the opposite of what it intends: the total column comes back empty. Excel
    and LibreOffice both write the value alongside the formula, so patching it
    in is what makes this fixture resemble a file a user would actually upload.
    """
    import re
    import zipfile

    with zipfile.ZipFile(path) as archive:
        names = archive.namelist()
        blobs = {name: archive.read(name) for name in names}

    sheet = next(name for name in names if name.endswith("sheet1.xml"))
    xml = blobs[sheet].decode("utf-8")
    xml, applied = re.subn(
        rf'(<c r="{cell}"[^>]*><f>[^<]+</f>)<v></v>',
        rf"\g<1><v>{value}</v>",
        xml,
    )
    if applied != 1:
        raise RuntimeError(
            f"Could not attach a cached value to {cell}; openpyxl's cell XML changed shape"
        )
    blobs[sheet] = xml.encode("utf-8")

    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
        for name in names:
            archive.writestr(name, blobs[name])


def make_spreadsheet_xlsx() -> None:
    """Two named worksheets, a merged header, and a formula with its result."""
    from openpyxl import Workbook

    workbook = Workbook()
    budget = workbook.active
    budget.title = "Бюджет курса"
    budget["A1"] = "Смета расходов на запуск"
    budget.merge_cells("A1:B1")
    budget["A2"] = "Статья"
    budget["B2"] = "Сумма"
    for row, (article, amount) in enumerate(
        (("Хостинг", 120), ("Лицензии", 340), ("Съёмка", 890)), start=3
    ):
        budget[f"A{row}"] = article
        budget[f"B{row}"] = amount
    budget["A6"] = "Итого"
    budget["B6"] = "=SUM(B3:B5)"

    staff = workbook.create_sheet("Преподаватели")
    staff["A1"] = "Роль"
    staff["B1"] = "Ставка"
    for row, (role, rate) in enumerate(
        (("Методист", 5600), ("Ассистент", 3100)), start=2
    ):
        staff[f"A{row}"] = role
        staff[f"B{row}"] = rate

    path = ROOT / "premium-spreadsheet.xlsx"
    workbook.save(path)
    _patch_xlsx_cached_formula(path, "B6", "1350")


def make_table_csv() -> None:
    """A plain delimited table; the family's whole structure is its columns."""
    (ROOT / "premium-table.csv").write_text(
        "Регион,Слушателей,Завершили\n"
        "Северный,120,84\n"
        "Южный,340,201\n"
        "Восточный,56,50\n",
        encoding="utf-8",
    )


def _odf_paragraph(text: str):
    from odf.text import P

    return P(text=text)


def make_opendocument_text() -> None:
    """ODT: headings, a nested-free list, and a table — ODF's defining shape."""
    from odf.opendocument import OpenDocumentText
    from odf.table import Table, TableCell, TableColumn, TableRow
    from odf.text import H, List, ListItem, P

    document = OpenDocumentText()
    document.text.addElement(H(outlinelevel=1, text="Регламент обучения"))
    document.text.addElement(
        P(text="Документ описывает порядок проведения корпоративного курса.")
    )
    document.text.addElement(H(outlinelevel=2, text="Этапы"))
    items = List()
    for item in (
        "Подготовка материалов",
        "Проведение занятий",
        "Итоговая аттестация",
    ):
        entry = ListItem()
        entry.addElement(_odf_paragraph(item))
        items.addElement(entry)
    document.text.addElement(items)
    document.text.addElement(H(outlinelevel=2, text="Контрольные показатели"))
    table = Table(name="Показатели")
    table.addElement(TableColumn(numbercolumnsrepeated=2))
    for row in (("Показатель", "Значение"), ("Явка", "92"), ("Аттестация", "78")):
        table_row = TableRow()
        for value in row:
            cell = TableCell(valuetype="string")
            cell.addElement(_odf_paragraph(value))
            table_row.addElement(cell)
        table.addElement(table_row)
    document.text.addElement(table)
    document.save(str(ROOT / "premium-opendocument.odt"))


def make_opendocument_sheet() -> None:
    """ODS: a named sheet whose name is the only thing tying rows to a subject."""
    from odf.opendocument import OpenDocumentSpreadsheet
    from odf.table import Table, TableCell, TableColumn, TableRow

    document = OpenDocumentSpreadsheet()
    table = Table(name="Нагрузка")
    table.addElement(TableColumn(numbercolumnsrepeated=2))
    for row in (("Модуль", "Часы"), ("Введение", "4"), ("Практика", "12")):
        table_row = TableRow()
        for value in row:
            cell = TableCell(valuetype="string")
            cell.addElement(_odf_paragraph(value))
            table_row.addElement(cell)
        table.addElement(table_row)
    document.spreadsheet.addElement(table)
    document.save(str(ROOT / "premium-opendocument.ods"))


def make_opendocument_slides() -> None:
    """ODP: three slides whose ORDER is the fact under test."""
    from odf.draw import Frame, Page, TextBox
    from odf.opendocument import OpenDocumentPresentation
    from odf.style import MasterPage, PageLayout, PageLayoutProperties

    document = OpenDocumentPresentation()
    layout = PageLayout(name="FixtureLayout")
    layout.addElement(
        PageLayoutProperties(margin="0cm", pagewidth="28cm", pageheight="21cm")
    )
    document.automaticstyles.addElement(layout)
    master = MasterPage(name="FixtureMaster", pagelayoutname=layout)
    document.masterstyles.addElement(master)

    slides = (
        "Первый слайд: цели программы",
        "Второй слайд: методика оценки",
        "Третий слайд: итоговые результаты",
    )
    for position, title in enumerate(slides):
        page = Page(masterpagename=master)
        document.presentation.addElement(page)
        frame = Frame(width="20cm", height="3cm", x="2cm", y=f"{2 + position}cm")
        box = TextBox()
        box.addElement(_odf_paragraph(title))
        frame.addElement(box)
        page.addElement(frame)
    document.save(str(ROOT / "premium-slides.odp"))


def make_epub() -> None:
    """EPUB: three chapters in spine order, each with a distinctive sentence."""
    from ebooklib import epub

    book = epub.EpubBook()
    book.set_identifier("mc2-premium-epub-fixture")
    book.set_title("Основы построения курса")
    book.set_language("ru")

    chapters = []
    for index, (title, body) in enumerate(
        (
            ("Глава первая. Введение", "Первая глава объясняет структуру программы."),
            ("Глава вторая. Практика", "Вторая глава содержит практические задания."),
            ("Глава третья. Аттестация", "Третья глава описывает критерии оценки."),
        ),
        start=1,
    ):
        chapter = epub.EpubHtml(title=title, file_name=f"chapter-{index}.xhtml", lang="ru")
        chapter.content = f"<h1>{title}</h1><p>{body}</p>"
        book.add_item(chapter)
        chapters.append(chapter)

    book.toc = tuple(chapters)
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav", *chapters]
    epub.write_epub(str(ROOT / "premium-book.epub"), book)


def make_latex() -> None:
    """LaTeX: a displayed equation and a verbatim listing, which is why the
    family is in scope at all — both arrive as typed items, not as prose."""
    (ROOT / "premium-paper.tex").write_text(
        r"""\documentclass{article}
\usepackage[utf8]{inputenc}
\begin{document}
\section{Модель усвоения}
Скорость усвоения материала описывается выражением
\begin{equation}
E = \frac{\alpha t}{1 + \beta t}
\end{equation}
\subsection{Реализация}
\begin{verbatim}
def retention(alpha, beta, t):
    return alpha * t / (1 + beta * t)
\end{verbatim}
\end{document}
""",
        encoding="utf-8",
    )


GENERATORS = {
    "enrichment-pdf": make_enrichment_pdf,
    "structured-docx": make_docx,
    "hierarchy-docx": make_hierarchy_docx,
    "numbered-sections-pdf": make_numbered_sections_pdf,
    "reading-order-pptx": make_pptx,
    "raster-ocr-pdf": make_raster_ocr_pdf,
    "vector-outline-pdf": make_vector_outline_pdf,
    "premium-spreadsheet-xlsx": make_spreadsheet_xlsx,
    "premium-table-csv": make_table_csv,
    "premium-opendocument-odt": make_opendocument_text,
    "premium-opendocument-ods": make_opendocument_sheet,
    "premium-slides-odp": make_opendocument_slides,
    "premium-book-epub": make_epub,
    "premium-paper-latex": make_latex,
}


if __name__ == "__main__":
    import sys

    ROOT.mkdir(parents=True, exist_ok=True)
    selected = sys.argv[1:] or list(GENERATORS)
    for name in selected:
        if name not in GENERATORS:
            raise SystemExit(f"Unknown fixture {name}; known: {', '.join(GENERATORS)}")
        GENERATORS[name]()
